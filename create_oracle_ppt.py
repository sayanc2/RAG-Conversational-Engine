#!/usr/bin/env python3
"""
ORACLE PowerPoint Presentation Generator
Creates a professional presentation about the RAG-based conversational engine
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import os

# Color scheme - Professional purple & gradient
COLOR_PRIMARY = RGBColor(102, 126, 234)      # Purple blue
COLOR_SECONDARY = RGBColor(118, 75, 162)     # Dark purple
COLOR_ACCENT = RGBColor(245, 158, 11)        # Orange accent
COLOR_TEXT_DARK = RGBColor(30, 27, 75)       # Dark text
COLOR_TEXT_LIGHT = RGBColor(243, 244, 246)   # Light text
COLOR_SUCCESS = RGBColor(16, 185, 129)       # Green
COLOR_WARN = RGBColor(239, 68, 68)           # Red/Warning

def add_title_slide(prs, title, subtitle=""):
    """Add a professional title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # Add gradient effect with shape
    left = top = Inches(0)
    shape = slide.shapes.add_shape(1, left, top, prs.slide_width, prs.slide_height)
    shape.fill.gradient()
    shape.fill.gradient_angle = 45.0
    shape.fill.gradient_stops[0].color.rgb = COLOR_PRIMARY
    shape.fill.gradient_stops[1].color.rgb = COLOR_SECONDARY
    shape.line.color.rgb = COLOR_PRIMARY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 210, 250)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_type="bullets"):
    """Add a content slide with title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_TEXT_LIGHT

    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT

    return slide


def add_bullet_content(slide, bullets, left=0.7, top=1.2):
    """Add bullet points to a slide."""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(8.8), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet['text']
        p.level = bullet.get('level', 0)
        p.font.size = Pt(18 - (bullet.get('level', 0) * 3))
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        if bullet.get('color'):
            p.font.color.rgb = bullet['color']


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add a two-column content slide."""
    slide = add_content_slide(prs, title)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(0.5))
    left_frame = left_box.text_frame
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    left_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(5))
    left_content = left_content_box.text_frame
    left_content.word_wrap = True

    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = left_content.paragraphs[0]
        else:
            p = left_content.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.5))
    right_frame = right_box.text_frame
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    right_content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.5), Inches(5))
    right_content = right_content_box.text_frame
    right_content.word_wrap = True

    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = right_content.paragraphs[0]
        else:
            p = right_content.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Divider
    divider = slide.shapes.add_shape(1, Inches(5.0), Inches(1.0), Inches(0.05), Inches(5.5))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_PRIMARY
    divider.line.color.rgb = COLOR_PRIMARY


def add_flow_diagram_slide(prs, title, steps):
    """Add a slide with a flow diagram."""
    slide = add_content_slide(prs, title)

    step_width = 1.8
    step_height = 0.8
    start_x = 0.8
    start_y = 1.5
    spacing = 0.3

    for i, step in enumerate(steps):
        x = start_x + (i * (step_width + spacing))

        if i < len(steps) - 1:
            # Draw arrow
            arrow = slide.shapes.add_connector(1, Inches(x + step_width), Inches(start_y + step_height/2),
                                              Inches(x + step_width + spacing - 0.15), Inches(start_y + step_height/2))
            arrow.line.color.rgb = COLOR_PRIMARY
            arrow.line.width = Pt(2)

        # Draw box
        box = slide.shapes.add_shape(1, Inches(x), Inches(start_y), Inches(step_width), Inches(step_height))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_PRIMARY
        box.line.color.rgb = COLOR_SECONDARY
        box.line.width = Pt(2)

        # Text
        text_frame = box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER


def add_table_slide(prs, title, table_data):
    """Add a slide with a table."""
    slide = add_content_slide(prs, title)

    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(4.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    col_width = width / cols
    for col_idx in range(cols):
        table_shape.columns[col_idx].width = int(col_width)

    # Fill table
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            # Header row styling
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_LIGHT
            else:
                cell.fill.solid()
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(243, 244, 246)
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def create_presentation():
    """Create the complete ORACLE presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ============ SLIDE 1: Title Slide ============
    add_title_slide(prs, "🔮 ORACLE",
                   "Orchestrated Retrieval and Conversational Logic Engine\nA RAG-based Multi-Agent Conversational System")

    # ============ SLIDE 2: Overview ============
    slide = add_content_slide(prs, "📊 System Overview")
    add_bullet_content(slide, [
        {'text': 'Advanced RAG-based conversational engine fusing structured + unstructured data', 'level': 0},
        {'text': 'SQL queries (employee data) + Vector search + Live web data (Tavily API)', 'level': 1},
        {'text': 'Multi-agent orchestration with specialized roles and handoff patterns', 'level': 0},
        {'text': 'Zero hallucination posture: validator agents verify ALL claims against sources', 'level': 0},
        {'text': 'Human-in-the-loop review for low-confidence answers (< 70% groundedness)', 'level': 0},
        {'text': 'Production-ready with comprehensive error handling and audit trails', 'level': 0},
    ])

    # ============ SLIDE 3: Problem Statement ============
    slide = add_content_slide(prs, "🎯 Problem Statement")
    add_bullet_content(slide, [
        {'text': 'Enterprise knowledge is fragmented across multiple data sources', 'level': 0},
        {'text': 'Employee databases (structured)', 'level': 1},
        {'text': 'Real-time news & weather (unstructured APIs)', 'level': 1},
        {'text': 'Traditional RAG systems suffer from hallucination and lack source grounding', 'level': 0},
        {'text': 'No unified interface for blended queries combining employee + contextual data', 'level': 0},
        {'text': 'Need for automated validation before any user-facing answer is delivered', 'level': 0},
        {'text': 'Human review required for edge cases and sensitive information combinations', 'level': 0},
    ])

    # ============ SLIDE 4: Solution Architecture ============
    slide = add_content_slide(prs, "🏗️ Solution Architecture")
    flow_steps = [
        "User Query",
        "Conductor\n(Triage)",
        "Specialist\nAgents",
        "Blended\nComposition",
        "Validation",
        "HITL Review",
        "Final Answer"
    ]
    add_flow_diagram_slide(prs, "🔄 Query Pipeline Flow", flow_steps)

    # ============ SLIDE 5: Multi-Agent System ============
    slide = add_content_slide(prs, "🤖 Multi-Agent System (5 Agents)")
    add_bullet_content(slide, [
        {'text': 'ORACLE Conductor (Triage & Composition)', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Routes queries, calls specialist tools, composes blended answers', 'level': 1},
        {'text': 'HERALD (Live News & Weather Specialist)', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Fetches from Tavily API, stores in Chroma, performs similarity search', 'level': 1},
        {'text': 'ARCHIVIST (SQL & Employee Knowledge Specialist)', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Queries SQLAlchemy database, semantic location mapping', 'level': 1},
        {'text': 'VALIDATOR (Groundedness Critic)', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Verifies factual claims, enforces zero hallucination via output guardrail', 'level': 1},
        {'text': 'SENTINEL (Security & Input/Output Guard)', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Prompt injection detection, PII protection, off-topic detection', 'level': 1},
    ])

    # ============ SLIDE 6: Manager Pattern (Conductor) ============
    slide = add_content_slide(prs, "🎭 Manager Pattern: ORACLE Conductor")
    add_bullet_content(slide, [
        {'text': 'Central orchestration agent that processes all user queries', 'level': 0},
        {'text': 'Implements intelligent triage routing:', 'level': 0},
        {'text': 'Pure weather/news query → full handoff to HERALD', 'level': 1},
        {'text': 'Pure employee query → full handoff to ARCHIVIST', 'level': 1},
        {'text': 'Blended query → calls both specialists as parallel tools & composes answer', 'level': 1},
        {'text': 'Handles PII-sensitive combinations requiring human review', 'level': 0},
        {'text': 'Primary model: Claude Sonnet 4-5 (with GPT-4o fallback)', 'level': 0},
        {'text': 'Attached guardrails: SENTINEL + VALIDATOR (mandatory security & groundedness)', 'level': 0},
    ])

    # ============ SLIDE 7: Tech Stack ============
    slide = add_content_slide(prs, "🛠️ Technology Stack")
    left_items = [
        "LLM & Orchestration",
        "• Claude Sonnet 4-5 (primary)",
        "• GPT-4o (fallback)",
        "• OpenAI Agents SDK",
        "",
        "Vector Store",
        "• ChromaDB (similarity search)",
        "• text-embedding-3-small",
    ]
    right_items = [
        "Database & Storage",
        "• SQLAlchemy (ORM)",
        "• SQLite + aiosqlite (async)",
        "• Redis (Phase 2)",
        "",
        "APIs & Frontend",
        "• Tavily (news + weather)",
        "• Streamlit (chat UI)",
    ]
    add_two_column_slide(prs, "📦 Tech Stack", "Backend", left_items, "Frontend & APIs", right_items)

    # ============ SLIDE 8: Agent Roster Details ============
    slide = add_content_slide(prs, "📋 Agent Roster Details")
    table_data = [
        ["Agent", "Role", "Model", "Guardrails", "SDK Features"],
        ["Conductor", "Triage & Composer", "Claude-Sonnet", "SENTINEL + VALIDATOR", "Handoffs + Tools"],
        ["HERALD", "Weather/News", "Claude-Sonnet", "None", "Full Agent + as_tool()"],
        ["ARCHIVIST", "Employee Data", "Claude-Sonnet", "None", "Full Agent + as_tool()"],
        ["VALIDATOR", "Groundedness Check", "GPT-4o", "Output Guardrail", "Claims verification"],
        ["SENTINEL", "Security Guard", "GPT-4o-mini", "Input + Output", "PII + Injection detection"],
    ]
    add_table_slide(prs, "📊 Agent Roster & Configuration", table_data)

    # ============ SLIDE 9: Semantic Join (Cross-Source Bridge) ============
    slide = add_content_slide(prs, "🔗 Semantic Join: Cross-Source Intelligence")
    add_bullet_content(slide, [
        {'text': 'Bridges employee database + live weather data without SQL JOIN', 'level': 0},
        {'text': 'HERALD fetches weather for location (e.g., "Austin, TX")', 'level': 0},
        {'text': 'Weather location embedded into vector space', 'level': 1},
        {'text': 'Queries ARCHIVIST\'s Chroma collection "employee_locations"', 'level': 0},
        {'text': 'Employee records matched by cosine similarity (distance < 0.30)', 'level': 1},
        {'text': 'Result: "Raghav Sharma works in Austin, TX. Weather there is 94°F, partly cloudy"', 'level': 0},
        {'text': 'Enables blended queries without explicit correlation schema', 'level': 0},
        {'text': 'Core architectural insight: embeddings solve impedance mismatch', 'level': 0},
    ])

    # ============ SLIDE 10: Zero Hallucination Posture ============
    slide = add_content_slide(prs, "🎯 Zero Hallucination Posture")
    add_bullet_content(slide, [
        {'text': 'VALIDATOR agent is MANDATORY output guardrail (never disabled)', 'level': 0},
        {'text': 'Process:', 'level': 0},
        {'text': 'Extract all factual claims from Conductor\'s answer', 'level': 1},
        {'text': 'Verify each claim against returned source chunks', 'level': 1},
        {'text': 'Score: % of grounded claims (0.0 - 1.0)', 'level': 1},
        {'text': 'Thresholds:', 'level': 0},
        {'text': 'Score ≥ 0.85: PASS → answer delivered', 'level': 1},
        {'text': 'Score 0.70-0.85: WARN → answer + warning logged', 'level': 1},
        {'text': 'Score < 0.70: FAIL → Human-in-the-loop activation', 'level': 1},
        {'text': 'If claim not in sources: explicitly state "data not available"', 'level': 0},
    ])

    # ============ SLIDE 11: Human-In-The-Loop (HITL) ============
    slide = add_content_slide(prs, "👥 Human-In-The-Loop (HITL) Framework")
    add_bullet_content(slide, [
        {'text': 'Triggered when validator score < 0.70 OR conductor flags PII sensitivity', 'level': 0},
        {'text': 'HITL Panel shows:', 'level': 0},
        {'text': 'Draft answer + groundedness score + ungrounded claims', 'level': 1},
        {'text': 'Human choices:', 'level': 0},
        {'text': '✅ Approve → answer published as-is', 'level': 1},
        {'text': '✏️ Edit & Approve → human-corrected version published', 'level': 1},
        {'text': '🔄 Regenerate → query re-run with rejection context', 'level': 1},
        {'text': 'All actions logged for audit trail & analytics', 'level': 0},
        {'text': 'Review time and human decisions tracked in session context', 'level': 0},
    ])

    # ============ SLIDE 12: HITL Hooks & Callbacks ============
    slide = add_content_slide(prs, "🪝 HITL Hooks & Lifecycle")
    add_bullet_content(slide, [
        {'text': 'Three-tier hook system:', 'level': 0},
        {'text': 'Agent-level: on_start, on_end (agent lifecycle)', 'level': 1},
        {'text': 'Run-level: on_agent_end with HITL detection', 'level': 1},
        {'text': 'HITL-specialized: on_hitl_triggered, on_hitl_approved, on_hitl_rejected', 'level': 1},
        {'text': 'Callbacks fire at each decision point:', 'level': 0},
        {'text': '[HITL TRIGGERED] → human review panel activated', 'level': 1},
        {'text': '[HITL APPROVED] or [HITL APPROVED (EDITED)] → answer published', 'level': 1},
        {'text': '[HITL REJECTED] → query re-run with new context', 'level': 1},
        {'text': 'Review metrics captured: duration, action, groundedness score', 'level': 0},
    ])

    # ============ SLIDE 13: Groundedness Checking ============
    slide = add_content_slide(prs, "✔️ Groundedness Validation Process")
    add_bullet_content(slide, [
        {'text': 'Two-phase validation:', 'level': 0},
        {'text': 'Phase 1: GPT-4o extracts factual claims from answer', 'level': 1},
        {'text': 'Phase 2: Match each claim against source document chunks', 'level': 1},
        {'text': 'Scoring:', 'level': 0},
        {'text': 'Grounded: claim explicitly supported by source', 'level': 1},
        {'text': 'Ungrounded: claim not found in any source', 'level': 1},
        {'text': 'Unknown: claim verifiable but sources incomplete', 'level': 1},
        {'text': 'Source types tracked: SQL rows, Chroma vectors, Tavily URLs', 'level': 0},
        {'text': 'Model mismatch (GPT-4o vs Claude): avoids confirmation bias', 'level': 0},
    ])

    # ============ SLIDE 14: Data Model & Schema ============
    slide = add_content_slide(prs, "📐 Data Models & Schema")
    table_data = [
        ["Component", "Type", "Details"],
        ["OracleSessionContext", "Pydantic", "session_id, user_id, conversation_history, HITL state"],
        ["ConductorResponse", "Pydantic", "answer, sources, confidence, query_type, hitl_required"],
        ["EmployeeRecord", "Pydantic", "employee_id, name, department, office_location"],
        ["GroundednessReport", "Pydantic", "score, claim_verifications, ungrounded_claims, verdict"],
        ["Employees DB", "SQLite", "500 rows, 10 cities, 8 departments (Raghav @ row 42, Austin TX)"],
        ["Chroma Collections", "Vector DB", "live_context (weather/news), employee_locations (embeddings)"],
    ]
    add_table_slide(prs, "📊 Data Models & Schema", table_data)

    # ============ SLIDE 15: Handoff Contracts ============
    slide = add_content_slide(prs, "📋 Handoff Contracts & Routing")
    add_bullet_content(slide, [
        {'text': 'Escalate to HERALD', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Trigger: Pure weather/news query', 'level': 1},
        {'text': 'Input type: HeraldHandoffInput (query + location_hint)', 'level': 1},
        {'text': 'Escalate to ARCHIVIST', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Trigger: Pure employee/HR query', 'level': 1},
        {'text': 'Input type: ArchivistHandoffInput (query + employee_hint + location_hint)', 'level': 1},
        {'text': 'Parallel Tool Calls (Blended Queries)', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'herald_as_tool("fetch_live_context") + archivist_as_tool("lookup_employee_data")', 'level': 1},
        {'text': 'Both execute simultaneously within same Conductor turn', 'level': 1},
    ])

    # ============ SLIDE 16: Development Workflow ============
    slide = add_content_slide(prs, "🔨 Development Workflow & Implementation Order")
    add_bullet_content(slide, [
        {'text': 'Phase 1A: Foundation (settings, DB, embeddings, Chroma)', 'level': 0},
        {'text': 'Phase 1B: Tools (SQL, Tavily, security, validation tools)', 'level': 0},
        {'text': 'Phase 1C: Agents (bottom-up: SENTINEL → VALIDATOR → specialists → Conductor)', 'level': 0},
        {'text': 'Phase 1D: Engine + UI (oracle_engine.py wrapper, Streamlit frontend)', 'level': 0},
        {'text': 'Phase 1E: Testing (unit, integration, E2E, acceptance test)', 'level': 0},
        {'text': 'Phase 2: MCP Extension (expose tools as MCP resources)', 'level': 0},
        {'text': 'Critical Path: Each phase built on previous layer', 'level': 0},
    ])

    # ============ SLIDE 17: Validation & Testing ============
    slide = add_content_slide(prs, "🧪 Validation & Testing Strategy")
    add_bullet_content(slide, [
        {'text': 'Canonical Demo Query: "What is the weather like where Raghav works?"', 'level': 0},
        {'text': 'Expected flow:', 'level': 0},
        {'text': 'SQL lookup → Raghav Sharma, Austin TX, Engineering', 'level': 1},
        {'text': 'Tavily fetch → Austin weather data', 'level': 1},
        {'text': 'Semantic match → EMP-0042 ↔ Austin location (distance < 0.30)', 'level': 1},
        {'text': 'Composition → blended answer with sources', 'level': 1},
        {'text': 'Validation → all claims verified (score: 0.97)', 'level': 1},
        {'text': 'SECURITY ✓ | GROUNDEDNESS ✓ | DELIVERY', 'level': 1},
        {'text': 'Test coverage: unit, integration, E2E, guardrails, edge cases', 'level': 0},
    ])

    # ============ SLIDE 18: API & Session Management ============
    slide = add_content_slide(prs, "📡 API & Session Management")
    add_bullet_content(slide, [
        {'text': 'Engine.run(user_query, ctx) → returns dict with:', 'level': 0},
        {'text': 'answer, response, error, hitl_triggered, security_blocked', 'level': 1},
        {'text': 'groundedness_score, hitl_metadata', 'level': 1},
        {'text': 'Session Persistence:', 'level': 0},
        {'text': 'Phase 1: SQLite session_memory table (cross-session in same user)', 'level': 1},
        {'text': 'Phase 2: Redis drop-in replacement (TTL: 24h)', 'level': 1},
        {'text': 'Conversation history captured in OracleSessionContext', 'level': 0},
        {'text': 'HITL decisions logged for audit trail & analytics', 'level': 0},
    ])

    # ============ SLIDE 19: Error Handling & Resilience ============
    slide = add_content_slide(prs, "🛡️ Error Handling & Resilience")
    add_bullet_content(slide, [
        {'text': 'InputGuardrailTripwireTriggered → friendly rejection message', 'level': 0},
        {'text': 'OutputGuardrailTripwireTriggered → HITL activation', 'level': 0},
        {'text': 'MaxTurnsExceeded → graceful degradation message', 'level': 0},
        {'text': 'TavilyAPIError → retry 3x via tenacity, fallback to cached Chroma', 'level': 0},
        {'text': 'SQLAlchemyError → log + return "data unavailable" message', 'level': 0},
        {'text': 'ChromaException → log + SQL-only fallback', 'level': 0},
        {'text': 'AnthropicAPIError → fallback to GPT-4o via RunConfig', 'level': 0},
        {'text': 'All errors logged with context for debugging', 'level': 0},
    ])

    # ============ SLIDE 20: Feature Highlights ============
    slide = add_content_slide(prs, "✨ Key Features & Innovations")
    add_bullet_content(slide, [
        {'text': 'Semantic Cross-Source Intelligence: embeddings bridge SQL ↔ APIs', 'level': 0},
        {'text': 'Agent-as-Tool Transformation: HERALD & ARCHIVIST work as full agents OR inline tools', 'level': 0},
        {'text': 'Zero Hallucination Enforcement: MANDATORY validator guardrail on all outputs', 'level': 0},
        {'text': 'Human-In-The-Loop with Audit Trail: track all review decisions & timings', 'level': 0},
        {'text': 'Multi-Model Strategy: Claude primary, GPT-4o for validation (avoids bias)', 'level': 0},
        {'text': 'Comprehensive Observability: hooks at agent/run/HITL levels with detailed logging', 'level': 0},
        {'text': 'Production-Ready: error handling, graceful degradation, session persistence', 'level': 0},
    ])

    # ============ SLIDE 21: Deployment & Operations ============
    slide = add_content_slide(prs, "🚀 Deployment & Operations")
    left_items = [
        "Environment Setup",
        "• Docker container support",
        "• Redis for distributed sessions",
        "• SQLite for local dev",
        "",
        "Configuration",
        "• .env file w/ API keys",
        "• Adjustable thresholds",
    ]
    right_items = [
        "Monitoring & Analytics",
        "• Detailed structured logs (JSONL)",
        "• HITL review metrics",
        "• Agent execution timing",
        "",
        "Scaling Considerations",
        "• Max 15 turns per query",
        "• Session TTL: 24h (configurable)",
    ]
    add_two_column_slide(prs, "🚀 Deployment & Operations", "Deployment", left_items, "Monitoring", right_items)

    # ============ SLIDE 22: Phase 2 Roadmap ============
    slide = add_content_slide(prs, "📈 Phase 2 Roadmap & Future Enhancements")
    add_bullet_content(slide, [
        {'text': 'MCP Server Extension', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Expose all tools as MCP resources for broader integration', 'level': 1},
        {'text': 'HITL Timeout & Auto-Escalation', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Auto-reject if human doesn\'t review within 5 minutes', 'level': 1},
        {'text': 'HITL Notifications', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Email/Slack alerts when HITL triggered', 'level': 1},
        {'text': 'Analytics Dashboard', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Track approval rates, review times, common ungrounded claims', 'level': 1},
        {'text': 'Redis Session Persistence', 'level': 0, 'color': COLOR_PRIMARY},
        {'text': 'Drop-in replacement for SQLite', 'level': 1},
    ])

    # ============ SLIDE 23: Questions & Demo ============
    slide = add_title_slide(prs, "Q&A", "Ready for your questions & live demo")

    # Save presentation
    output_path = Path(__file__).parent / "ORACLE_Presentation.pptx"
    prs.save(str(output_path))
    print("[OK] Presentation created: {}".format(output_path))
    print("   - 23 slides")
    print("   - Professional design with gradient headers")
    print("   - Complete coverage of all topics")
    print("   - Ready for large-audience demonstration")


if __name__ == "__main__":
    create_presentation()
