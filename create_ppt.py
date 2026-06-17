"""
Create a PowerPoint presentation on the ORACLE codebase.
Covers architecture, concepts, components, and implementation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PURPLE = RGBColor(102, 126, 234)  # #667eea
DARK_PURPLE = RGBColor(118, 75, 162)  # #764ba2
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(30, 27, 75)
LIGHT_TEXT = RGBColor(148, 163, 184)

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 12, 41)  # dark bg

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_TEXT

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PURPLE
    title_shape.line.color.rgb = PURPLE

    # Title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add a two-column slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PURPLE
    title_shape.line.color.rgb = PURPLE

    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(0.4))
    lh_frame = left_header.text_frame
    lh_frame.text = left_title
    lh_frame.paragraphs[0].font.size = Pt(20)
    lh_frame.paragraphs[0].font.bold = True
    lh_frame.paragraphs[0].font.color.rgb = DARK_PURPLE

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4.2), Inches(4.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, point in enumerate(left_points):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column
    right_header = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4), Inches(0.4))
    rh_frame = right_header.text_frame
    rh_frame.text = right_title
    rh_frame.paragraphs[0].font.size = Pt(20)
    rh_frame.paragraphs[0].font.bold = True
    rh_frame.paragraphs[0].font.color.rgb = DARK_PURPLE

    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(2.3), Inches(4.2), Inches(4.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, point in enumerate(right_points):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# ============================================================================
# SLIDES
# ============================================================================

# Slide 1: Title
add_title_slide(prs, "🔮 ORACLE", "Orchestrated Retrieval and Conversational Logic Engine\nA Multi-Agent RAG System with Human-In-The-Loop Validation")

# Slide 2: System Overview
add_content_slide(prs, "System Overview", [
    "🎯 Purpose: Blend structured (SQL) + unstructured (vector/web) data retrieval",
    "🤖 5 Specialized Agents working in orchestrated harmony",
    "📊 RAG Pattern: Employee database + Live weather/news + Vector embeddings",
    "🛡️ Safety-first: Dual guardrails (security + groundedness)",
    "👥 HITL: Human review when AI confidence drops below threshold"
])

# Slide 3: The 5 Agents
add_two_column_slide(prs, "The 5-Agent Architecture",
    "Specialist Agents",
    [
        "🌤️ HERALD: Live weather & news via Tavily API",
        "👤 ARCHIVIST: SQL employee queries + semantic location mapping",
        "✅ VALIDATOR: Groundedness critic (GPT-4o, independent perspective)",
        "🛡️ SENTINEL: Security guardrail (prompt injection, PII detection)"
    ],
    "Orchestration",
    [
        "🎭 CONDUCTOR: Triage + composition brain",
        "• Routes intent to specialists",
        "• Calls tools for blended queries",
        "• Composes final answer",
        "• Compiles sources for HITL"
    ]
)

# Slide 4: Data Model
add_content_slide(prs, "Data Architecture", [
    "🗄️ SQLite + aiosqlite: 500 employee records (Raghav Sharma @ EMP-0042, Austin TX)",
    "🔍 ChromaDB: Two collections",
    "    • live_context: Weather/news embeddings from Tavily (ephemeral)",
    "    • employee_locations: 500 embedded office locations (persistent)",
    "🌐 Tavily API: Real-time weather & news for any location",
    "📡 Semantic Bridge: Cosine similarity <0.30 matches weather location to employees"
])

# Slide 5: Request Flow
add_content_slide(prs, "Canonical Request Flow: 'What is the weather where Raghav works?'", [
    "1️⃣ SENTINEL input guardrail checks for prompt injection/PII extraction",
    "2️⃣ CONDUCTOR detects blended query intent (employee + weather)",
    "3️⃣ ARCHIVIST tool finds Raghav → Austin, TX, Engineering",
    "4️⃣ HERALD tool fetches Austin weather from Tavily → embeds into live_context",
    "5️⃣ Semantic match: Austin weather ↔ employee_locations via Chroma",
    "6️⃣ CONDUCTOR composes: 'Raghav works in Austin. Weather: 94°F, Partly Cloudy.'",
    "7️⃣ VALIDATOR critic verifies all claims against sources → score 0.97 (pass)",
    "8️⃣ SENTINEL output guardrail checks for PII leak",
    "9️⃣ Answer + sources delivered to UI"
])

# Slide 6: HITL (Human-In-The-Loop)
add_content_slide(prs, "Human-In-The-Loop (HITL) Flow", [
    "🎯 Trigger: Groundedness score < 0.70 OR sensitive PII combination detected",
    "📋 Panel shows:",
    "    • Draft answer with ungrounded claims highlighted",
    "    • Groundedness bar (🟢 pass / 🟡 warn / 🔴 fail)",
    "    • Editable text box for human refinement",
    "👥 Human chooses:",
    "    ✅ Approve: Publish as-is",
    "    ✏️ Edit & Approve: Publish human-edited version",
    "    🔄 Regenerate: Re-run with error feedback"
])

# Slide 7: Tech Stack
add_two_column_slide(prs, "Technology Stack",
    "Core Dependencies",
    [
        "🤖 OpenAI Agents SDK (orchestration)",
        "🧠 Claude Sonnet 4.5 (primary model)",
        "🔀 GPT-4o (validation/guardrails)",
        "🎨 Streamlit (frontend)",
        "🔍 ChromaDB (vector store)",
        "💾 SQLAlchemy + aiosqlite (ORM)"
    ],
    "APIs & Services",
    [
        "🌐 Tavily API (weather/news)",
        "🔑 OpenAI Embeddings (text-embedding-3-small)",
        "⚙️ OpenAI API (GPT-4o)",
        "🧭 Anthropic API (Claude)",
        "📝 Python 3.14 (async/await)",
        "🔐 Environment variables (.env)"
    ]
)

# Slide 8: Guardrails & Safety
add_content_slide(prs, "Safety Architecture: Dual Guardrails", [
    "🛡️ SENTINEL (GPT-4o-mini, fast & cheap):",
    "    INPUT: Detects prompt injection, off-topic, PII extraction",
    "    OUTPUT: Detects PII leaks, bulk data exports",
    "",
    "✅ VALIDATOR (GPT-4o, independent model):",
    "    Verifies every factual claim against retrieved sources",
    "    Score calculation: grounded_claims / total_claims",
    "    Thresholds: ≥0.85 pass | 0.70-0.85 warn | <0.70 HITL trigger",
    "",
    "🔒 Both run in parallel (fast) + post-generation (accurate)"
])

# Slide 9: Streamlit UI Features
add_content_slide(prs, "Streamlit Frontend Features", [
    "💬 Chat Interface: Demo query buttons + message input",
    "📊 Left Sidebar: Session metrics, groundedness bar, source explorer",
    "🎨 Custom CSS: Gradient purple theme, responsive cards",
    "❗ HITL Panel: Auto-appears on low groundedness with approve/edit/regenerate",
    "🏷️ Badges: Confidence levels, query type, source type indicators",
    "💡 Follow-up Suggestions: 2-3 related queries recommended",
    "📝 Session Persistence: Conversation history saved to SQLite"
])

# Slide 10: Testing
add_content_slide(prs, "Test Coverage (18/18 Passing)", [
    "🧪 test_tools.py: Unit tests for all 6 function tools",
    "    • Embedding, Chroma, SQL, Tavily, Security, Validation",
    "",
    "🛡️ test_guardrails.py: Guardrail tripwire behavior",
    "    • SENTINEL blocks injection, VALIDATOR triggers HITL",
    "",
    "🔀 test_blended_query.py: Canonical acceptance test",
    "    • 'What is the weather where Raghav works?'",
    "    • Tests full orchestration, HITL, state management",
    "",
    "✅ All critical paths covered. Ready for production."
])

# Slide 11: Deployment
add_content_slide(prs, "How to Run ORACLE", [
    "📋 Prerequisites:",
    "    ANTHROPIC_API_KEY, OPENAI_API_KEY, TAVILY_API_KEY in oracle/.env",
    "",
    "🌱 Seed Database (run once):",
    "    PROTOCOL_BUFFERS_PYTHON_IMPLEMENTATION=python python -m oracle.seed_db",
    "",
    "🚀 Launch Application:",
    "    PROTOCOL_BUFFERS_PYTHON_IMPLEMENTATION=python python -m streamlit run oracle/app.py",
    "",
    "🌐 Access at http://localhost:8501"
])

# Slide 12: Key Concepts Summary
add_content_slide(prs, "Key Architectural Concepts", [
    "✨ Agent Duality: Agents work as FULL agents (via handoff) OR inline TOOLS (via as_tool())",
    "🔗 Semantic Bridge: Chroma similarity search connects weather ↔ employees across modalities",
    "🎯 Input Filtering: Conversation history scrubbed per agent (HERALD strips employee turns, etc.)",
    "🏆 Parallel Execution: Input guardrail runs in parallel with agent startup (low latency)",
    "💾 Session State: Pydantic context object persists across turns + survives SQLite restarts",
    "🎨 Cascade Safety: 3-layer validation (input guard → validator → output guard)"
])

# Slide 13: Lessons Learned
add_content_slide(prs, "Development Insights", [
    "🐍 Python 3.14: Requires PROTOCOL_BUFFERS_PYTHON_IMPLEMENTATION=python workaround",
    "🔧 function_tool: No __wrapped__ attribute; define raw _fn + wrap separately for testing",
    "📦 Lazy Imports: oracle/agents/__init__.py uses __getattr__ to avoid circular imports",
    "🎯 Strict Schema: dict[str, Any] params need strict_mode=False to pass SDK validation",
    "🚀 Optimization: Lazy-load engine in Streamlit callbacks to avoid startup blocking"
])

# Slide 14: Future Roadmap
add_content_slide(prs, "Phase 2: MCP Extension", [
    "🌐 MCP Server: Expose ORACLE tools as MCP resources for any LLM client",
    "🔌 Tool Wrapping: Convert function_tools → MCP resources (zero architectural change)",
    "🧠 Multi-Model: MCP allows Claude, GPT, open-source models to call ORACLE tools",
    "📡 Distributed: Run agents in separate MCP server instances if needed",
    "⚡ Performance: Caching + vectorization for 10k+ employee datasets"
])

# Slide 15: Closing
add_title_slide(prs, "Questions?", "ORACLE: Where Retrieval Meets Orchestration\n18/18 Tests Passing • Production Ready")

# Save
prs.save('ORACLE_Presentation.pptx')
print("PowerPoint created: ORACLE_Presentation.pptx")
